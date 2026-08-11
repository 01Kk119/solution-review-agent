# -*- coding: utf-8 -*-
"""PPTX 抽取：slide 标题、正文、表格、图片、speaker notes，均带 slide 编号。"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import (KIND_IMAGE, KIND_NOTE, KIND_TABLE, KIND_TEXT,
                   FileExtraction, Unit, rows_to_markdown_table, sanitize_name)


def _walk_shapes(shapes):
    """递归展开组合形状。"""
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(sh.shapes)
        else:
            yield sh


def extract_pptx(path: Path, fx: FileExtraction, assets_dir: Path, cfg: dict) -> None:
    min_bytes = cfg.get("min_image_bytes", 3000)
    stem = sanitize_name(path.stem, 40)
    prs = Presentation(str(path))

    for sno, slide in enumerate(prs.slides, start=1):
        locator = f"slide {sno}"
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
        except Exception:  # noqa: BLE001
            pass

        texts, img_idx = [], 0
        for sh in _walk_shapes(slide.shapes):
            # 文本
            if sh.has_text_frame:
                t = "\n".join(p.text for p in sh.text_frame.paragraphs if p.text.strip())
                if t.strip() and t.strip() != title:
                    texts.append(t.strip())
            # 表格
            if getattr(sh, "has_table", False) and sh.has_table:
                rows = [[c.text for c in r.cells] for r in sh.table.rows]
                md = rows_to_markdown_table(rows)
                if md:
                    fx.add(Unit(locator=f"{locator} table", kind=KIND_TABLE,
                                text=md, context=f"slide 标题：{title}"))
            # 图片
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = sh.image.blob
                    if len(blob) < min_bytes:
                        continue
                    img_idx += 1
                    ext = sh.image.ext or "png"
                    name = f"{fx.file_id}_{stem}_s{sno}_img{img_idx}.{ext}"
                    out = assets_dir / "images" / name
                    out.parent.mkdir(parents=True, exist_ok=True)
                    out.write_bytes(blob)
                    fx.add(Unit(locator=f"{locator} image {img_idx}", kind=KIND_IMAGE,
                                asset_path=f"assets/images/{name}",
                                needs_visual_reading=True,
                                context=f"slide 标题：{title}"))
                except Exception as e:  # noqa: BLE001
                    fx.notes.append(f"{locator}: 图片导出失败（{e}）")

        body = ("\n".join(texts)).strip()
        combined = (f"【标题】{title}\n{body}" if title else body).strip()
        if combined:
            fx.add(Unit(locator=locator, kind=KIND_TEXT, text=combined))

        # speaker notes
        try:
            if slide.has_notes_slide:
                nt = slide.notes_slide.notes_text_frame.text.strip()
                if nt:
                    fx.add(Unit(locator=f"{locator} notes", kind=KIND_NOTE, text=nt,
                                context=f"slide 标题：{title}"))
        except Exception as e:  # noqa: BLE001
            fx.notes.append(f"{locator}: speaker notes 读取失败（{e}）")
