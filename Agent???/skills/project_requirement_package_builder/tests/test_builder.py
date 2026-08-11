# -*- coding: utf-8 -*-
"""project_requirement_package_builder 最小测试。

覆盖点（对应需求）：
  1. 能读取输入目录            → test_extract_end_to_end
  2. 能识别不同文件类型        → test_extract_end_to_end
  3. 能生成 Markdown（草稿）   → test_scaffold
  4. 能生成 metadata.json      → test_validate_full_package
  5. 能生成 evidence_index.json→ test_validate_full_package
  6. 缺失信息写入 checklist    → test_scaffold
  7. 单文件解析失败不崩溃      → test_extract_end_to_end（broken.pdf）
  8. 解析失败写入处理日志      → test_processing_log

运行：cd skills/project_requirement_package_builder && python3 -m pytest tests/ -q
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest

SKILL_DIR = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(SKILL_DIR / "src"))

import index as cli  # noqa: E402


# ---------------------------------------------------------------- fixtures

def _make_png_bytes(w=120, h=120) -> bytes:
    import fitz
    pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, w, h), False)
    pix.clear_with(90)
    return pix.tobytes("png")


def _build_input_dir(root: Path) -> Path:
    """构造覆盖全部类型的 mock 项目资料目录。"""
    inp = root / "input_project"
    (inp / "06_transcripts").mkdir(parents=True)

    # pdf：第 1 页有文字，第 2 页纯图片（触发快照 + 视觉读取标记）
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 100), "Customer needs 3 forklifts for pallet handling.")
    page2 = doc.new_page()
    page2.insert_image(fitz.Rect(50, 50, 400, 400), stream=_make_png_bytes(400, 400))
    doc.save(inp / "solution.pdf")
    doc.close()

    # pptx：标题 + 正文 + 表格 + 备注
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "取货场景"
    slide.placeholders[1].text = "从 A 区地面取 1200x1000 托盘"
    rows, cols = 2, 2
    tb = slide.shapes.add_table(rows, cols, Inches(1), Inches(3),
                                Inches(4), Inches(1)).table
    tb.cell(0, 0).text = "车型"
    tb.cell(0, 1).text = "数量"
    tb.cell(1, 0).text = "VNE40-66"
    tb.cell(1, 1).text = "8"
    slide.notes_slide.notes_text_frame.text = "客户口头确认过数量"
    prs.save(inp / "plan.pptx")

    # xlsx：含合并单元格
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "基本信息"
    ws["A1"] = "项目"
    ws["B1"] = "N068"
    ws["A2"] = "车型"
    ws["B2"] = "VNE40-66"
    ws.merge_cells("A3:B3")
    ws["A3"] = "备注：分两期交付"
    wb.save(inp / "info.xlsx")

    # docx：标题 + 段落 + 表格
    import docx
    d = docx.Document()
    d.add_heading("Project Information", level=1)
    d.add_paragraph("Deliver AGV system to Napoleon site.")
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "Item"
    t.cell(0, 1).text = "Qty"
    t.cell(1, 0).text = "AGV"
    t.cell(1, 1).text = "8"
    d.save(inp / "spec.docx")

    # md / txt（转写带时间戳） / 独立图片
    (inp / "summary.md").write_text("# 会议纪要\n\n## 结论\n货叉宽度 125mm。\n",
                                    encoding="utf-8")
    (inp / "06_transcripts" / "transcript.txt").write_text(
        "00:00:01 主持人：开始。\n00:12:35 工程师：通道宽度 4.6 米。\n",
        encoding="utf-8")
    (inp / "site_photo.png").write_bytes(_make_png_bytes())

    # 损坏文件：伪 PDF（测试单文件失败不中断）
    (inp / "broken.pdf").write_bytes(b"this is not a real pdf content x" * 4)

    # 视频：仅登记
    (inp / "recording.mp4").write_bytes(b"\x00" * 2048)
    return inp


@pytest.fixture(scope="module")
def workspace(tmp_path_factory):
    root = tmp_path_factory.mktemp("prpb")
    inp = _build_input_dir(root)
    out = root / "output_project"
    rc = cli.main(["extract", "--input", str(inp), "--output", str(out),
                   "--project-name", "测试项目"])
    return {"root": root, "input": inp, "output": out, "rc": rc}


def _manifest(ws) -> dict:
    return json.loads((ws["output"] / "extracted" / "manifest.json")
                      .read_text(encoding="utf-8"))


# ---------------------------------------------------------------- tests

def test_extract_end_to_end(workspace):
    assert workspace["rc"] == 0
    m = _manifest(workspace)
    by_name = {f["source_file"]: f for f in m["files"]}

    # 类型识别
    assert by_name["solution.pdf"]["source_type"] == "pdf"
    assert by_name["plan.pptx"]["source_type"] == "pptx"
    assert by_name["info.xlsx"]["source_type"] == "xlsx"
    assert by_name["spec.docx"]["source_type"] == "docx"
    assert by_name["summary.md"]["source_type"] == "md"
    assert by_name["06_transcripts/transcript.txt"]["source_type"] == "txt"
    assert by_name["site_photo.png"]["source_type"] == "image"

    # 损坏文件失败但整体流程不崩溃，其余文件照常成功
    assert by_name["broken.pdf"]["parse_status"] == "failed"
    assert by_name["broken.pdf"]["error"]
    assert by_name["solution.pdf"]["parse_status"] == "ok"

    # 视频仅登记
    assert by_name["recording.mp4"]["parse_status"] == "skipped_media"

    # 每个文件都有对应的 extracted json
    for f in m["files"]:
        assert (workspace["output"] / f["extracted_json"]).exists()


def test_pdf_snapshot_and_visual_flag(workspace):
    m = _manifest(workspace)
    pdf = next(f for f in m["files"] if f["source_file"] == "solution.pdf")
    assert pdf["stats"]["needs_visual_reading"] >= 1  # 纯图片页触发
    data = json.loads((workspace["output"] / pdf["extracted_json"])
                      .read_text(encoding="utf-8"))
    kinds = {u["kind"] for u in data["units"]}
    assert "text" in kinds
    assert kinds & {"snapshot", "image"}
    for u in data["units"]:
        if u.get("asset_path"):
            assert (workspace["output"] / u["asset_path"]).exists()


def test_table_extraction(workspace):
    m = _manifest(workspace)
    for name in ("info.xlsx", "spec.docx", "plan.pptx"):
        f = next(x for x in m["files"] if x["source_file"] == name)
        assert f["stats"]["table_units"] >= 1, f"{name} 应至少抽出一个表格"
    xlsx = next(f for f in m["files"] if f["source_file"] == "info.xlsx")
    data = json.loads((workspace["output"] / xlsx["extracted_json"])
                      .read_text(encoding="utf-8"))
    table = next(u for u in data["units"] if u["kind"] == "table")
    assert "VNE40-66" in table["text"]
    assert "合并" in table["text"]  # 合并单元格被标注


def test_transcript_timestamp_locator(workspace):
    m = _manifest(workspace)
    t = next(f for f in m["files"] if f["source_file"].endswith("transcript.txt"))
    data = json.loads((workspace["output"] / t["extracted_json"])
                      .read_text(encoding="utf-8"))
    assert any("时间戳" in u["locator"] for u in data["units"])


def test_scaffold(workspace):
    rc = cli.main(["scaffold", "--output", str(workspace["output"]),
                   "--project-name", "测试项目"])
    assert rc == 0
    draft = (workspace["output"] / "project_requirement_package.draft.md")
    assert draft.exists()
    text = draft.read_text(encoding="utf-8")
    # 0.2 输入资料范围覆盖所有文件；模板占位符已替换
    for name in ("solution.pdf", "broken.pdf", "recording.mp4"):
        assert name in text
    assert "{{project_name}}" not in text
    assert (workspace["output"] / "missing_info_checklist.draft.md").exists()


def test_processing_log(workspace):
    text = (workspace["output"] / "project_requirement_package.draft.md") \
        .read_text(encoding="utf-8")
    log = text.split("## 13. AI 处理日志")[1]
    assert "解析失败文件" in log and "broken.pdf" in log
    assert "recording.mp4" in log  # 跳过的媒体文件也在日志中


def _write_minimal_final_package(out: Path):
    """构造一份最小但合法的最终输出，验证 render + validate。"""
    sections = "\n\n".join(
        f"## {i}. {t}\n\n（略）" for i, t in enumerate(
            ["文档说明", "项目基础信息", "业务流程与作业流", "车辆与产品形态",
             "搬运对象与载具信息", "取货 / 放货 / 识别场景", "导航与现场环境",
             "调度、接口与软件模块", "硬件与安全相关信息", "项目需求清单",
             "待确认问题清单", "资料侧初步风险提示", "原始资料索引", "AI 处理日志"]))
    md = ("# 项目订单原始资料包\n\n" + sections +
          "\n\n货叉宽度 125mm（来源：summary.md；证据 E001）\n")
    (out / "project_requirement_package.md").write_text(md, encoding="utf-8")
    (out / "metadata.json").write_text(json.dumps({
        "project_name": "测试项目", "customer_name": "未在资料中明确提供",
        "region": "未在资料中明确提供", "industry": "未在资料中明确提供",
        "vehicle_models": ["VNE40-66"], "vehicle_count": "8",
        "source_files": [{"file": "summary.md", "type": "md", "parse_status": "ok"}],
        "generated_at": "2026-07-08T00:00:00+08:00",
        "skill_version": "0.1.0", "language": "zh-CN",
    }, ensure_ascii=False), encoding="utf-8")
    (out / "evidence_index.json").write_text(json.dumps([{
        "id": "E001", "source_file": "summary.md", "source_type": "md",
        "page_or_slide": "章节「结论」", "extracted_text": "货叉宽度 125mm。",
        "used_in_sections": ["9. 项目需求清单"], "confidence": "原文明确",
    }], ensure_ascii=False), encoding="utf-8")
    (out / "missing_info_checklist.md").write_text(
        "# 缺失信息清单\n\n| # | 缺失信息 |\n|---|---|\n| 1 | 客户名称 |\n",
        encoding="utf-8")


def test_validate_full_package(workspace):
    out = workspace["output"]
    _write_minimal_final_package(out)
    assert cli.main(["render", "--output", str(out)]) == 0
    html = (out / "project_requirement_package.html").read_text(encoding="utf-8")
    assert "<table" in html or "项目订单原始资料包" in html
    assert cli.main(["validate", "--output", str(out)]) == 0


def test_validate_catches_problems(workspace, tmp_path):
    out = tmp_path / "bad_output"
    out.mkdir()
    # 缺文件
    assert cli.main(["validate", "--output", str(out)]) == 1
    # 引用不存在的证据编号
    _write_minimal_final_package(out)
    md = out / "project_requirement_package.md"
    md.write_text(md.read_text(encoding="utf-8") + "\n幽灵证据（证据 E999）\n",
                  encoding="utf-8")
    assert cli.main(["validate", "--output", str(out)]) == 1
